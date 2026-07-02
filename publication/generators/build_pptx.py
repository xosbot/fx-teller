"""
FX Teller Strategy Deck Generator
Builds an executive PowerPoint deck from the knowledge base.

Output: publication/build/Deck.pptx
"""
import os
import sys
from pathlib import Path
from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

REPO_ROOT = Path("/mnt/d/FX Teller")
PUB_ROOT = REPO_ROOT / "publication"
OUT_DIR = PUB_ROOT / "build"
PPTX_PATH = OUT_DIR / "Deck.pptx"

sys.path.insert(0, str(PUB_ROOT / "generators"))
sys.path.insert(0, str(PUB_ROOT))
from content_loader import build_index
from diagrams.library import render_diagram

# Brand colors
NAVY = RGBColor(0x1B, 0x1B, 0x2F)
GOLD = RGBColor(0xC5, 0xA5, 0x5A)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SAGE = RGBColor(0x7A, 0x9E, 0x7E)
SLATE = RGBColor(0x4A, 0x4A, 0x6A)
CORAL = RGBColor(0xC8, 0x7A, 0x6A)
LIGHT = RGBColor(0xF8, 0xF6, 0xF2)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6E, 0x6E, 0x6E)
BORDER = RGBColor(0xE0, 0xDD, 0xD5)


def add_bg(slide, color):
    """Add a colored background to a slide."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide.part.package.presentation_part.presentation.slide_width, slide.part.package.presentation_part.presentation.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=NAVY,
             font="Inter", align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=14, color=NAVY, font="Inter"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def add_rule(slide, top, color=GOLD, width_inches=2, left_inches=4.5):
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_inches), Inches(top), Inches(width_inches), Emu(25400))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()
    return rule


def add_image_from_svg(slide, svg_content, left, top, width, height):
    """Render an SVG to PNG and insert as image."""
    import cairosvg
    png_data = cairosvg.svg2png(bytestring=svg_content.encode("utf-8"),
                                  output_width=int(width.emu / 9525 * 96),
                                  output_height=int(height.emu / 9525 * 96))
    slide.shapes.add_picture(BytesIO(png_data), left, top, width=width, height=height)


def add_diagram(slide, diagram_name, left, top, width, height):
    """Add a SVG diagram to the slide."""
    svg = render_diagram(diagram_name)
    try:
        import cairosvg
        png_data = cairosvg.svg2png(bytestring=svg.encode("utf-8"),
                                      output_width=1200, output_height=1200)
        slide.shapes.add_picture(BytesIO(png_data), left, top, width=width, height=height)
        return True
    except ImportError:
        # Fallback: add text placeholder
        add_text(slide, left, top, width, height, f"[{diagram_name}]", size=14, color=MUTED, align=PP_ALIGN.CENTER)
        return False


def add_footer(slide, page_num, total_pages):
    """Add footer to slide."""
    sw = Inches(13.333)
    sh = Inches(0.4)
    # Footer line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.1), Inches(12.133), Emu(7620))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    # Left text
    add_text(slide, Inches(0.6), Inches(7.2), Inches(6), Inches(0.3),
             "FX Teller  |  Strategic Blueprint", size=8, color=MUTED)
    # Right page number
    add_text(slide, Inches(10.5), Inches(7.2), Inches(2.2), Inches(0.3),
             f"{page_num} / {total_pages}", size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def make_slide(prs, color_bg=WHITE):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    if color_bg != WHITE:
        add_bg(slide, color_bg)
    return slide


def main():
    print("Loading content...")
    index = build_index()
    total_words = index['total_words']
    total_docs = index['total_docs']
    sections = index['sections']

    # 16:9 widescreen
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Compute total page count first
    slides = []

    # Slide 1: Cover
    s = make_slide(prs, color_bg=NAVY)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.5),
             "FX TELLER", size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rule1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(3.0), Inches(2.5), Emu(12700))
    rule1.fill.solid(); rule1.fill.fore_color.rgb = GOLD; rule1.line.fill.background()
    add_text(s, Inches(0.5), Inches(3.2), Inches(12.333), Inches(1.0),
             "Strategic", size=44, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.9),
             "Blueprint", size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.4),
             "Founder Edition", size=14, color=WHITE, align=PP_ALIGN.CENTER)
    rule2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(5.5), Inches(1.3), Emu(12700))
    rule2.fill.solid(); rule2.fill.fore_color.rgb = GOLD; rule2.line.fill.background()
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.333), Inches(0.3),
             "Version 1.0  |  July 2026", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.05), Inches(12.333), Inches(0.3),
             "Prepared by the Executive Strategy Office", size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.3),
             "CONFIDENTIAL  |  For Founder, Board, and Investor Distribution", size=8, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
    slides.append(s)

    # Slide 2: Letter to the Founder (excerpt)
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PROLOGUE", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
             "Letter to the Founder", size=28, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(4.5),
             "This document is the Master Strategic Blueprint for FX Teller. It is the company in writing. "
             "It has been produced by the full Executive Strategy Office, drawing on every strategic document "
             "in the knowledge base, and stress-tested by the most-sceptical reader the Office could find. "
             "It is the document you will hand to the next executive, the next investor, the next regulator, the next Host.\n\n"
             "What has been built in this knowledge base is unusual. Twenty-two strategic decisions recorded in the Decision Register, "
             "with named owners and review cadences. Three governance frameworks: a 10-specialist Strategy Office with explicit vetoes, "
             "a 6-rhythm operational cadence, a documentation discipline that survives a 90-day Founder absence.\n\n"
             "A legal stack: SEBI opinion, Member agreement suite, data protection posture, audit trails. "
             "A financial model with Conservative, Expected, and Optimistic scenarios. A 6-month execution plan with week-by-week milestones. "
             "A 24-month strategic roadmap with quarterly sequencing.",
             size=12, color=DARK)
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "— The Executive Strategy Office", size=10, color=MUTED, italic=True, align=PP_ALIGN.RIGHT)
    add_footer(s, 2, 18)
    slides.append(s)

    # Slide 3: Executive Summary
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "EXECUTIVE SUMMARY", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
             "The Company, In One Sentence", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_text(s, Inches(0.6), Inches(2.1), Inches(12), Inches(1.5),
             "FX Teller is a members-only, audio-first Trading Floor for disciplined retail traders. "
             "A members' club whose principal amenity is a live, audio-only Trading Floor where a Host narrates the market in real time. "
             "The brand is the philosophy: discipline over engagement, process over prediction, calm over urgency.",
             size=14, color=DARK, italic=True)
    # Stat boxes row
    stats = [
        ("30–50", "Members Y1"),
        ("150–200", "Members Y2"),
        ("350–500", "Members Y3"),
        ("₹25L", "MRR at M24"),
        ("75%", "Gross Margin"),
    ]
    box_w = Inches(2.4)
    box_h = Inches(1.0)
    box_y = Inches(4.0)
    start_x = Inches(0.6)
    gap = Inches(0.05)
    for i, (val, label) in enumerate(stats):
        x = start_x + i * (box_w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, box_y, box_w, box_h)
        box.fill.solid(); box.fill.fore_color.rgb = CREAM; box.line.fill.background()
        add_text(s, x, box_y + Inches(0.1), box_w, Inches(0.5), val, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x, box_y + Inches(0.65), box_w, Inches(0.3), label, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5),
             "Capital posture: bootstrap Year 1 (₹1.2–1.5 Cr Founder reserve), Year 2 strategic raise ₹15–30 Cr on philosophy-defending terms. "
             "Team: 4 in Year 1, 6 by Q4, 8 by Year 2. House of Traders: scoped Q5, built Q7, opened Q13–Q14. "
             "AI Companion: production by Q11–Q12. The plan is a Founder-protection plan disguised as a launch plan.",
             size=12, color=DARK)
    add_footer(s, 3, 18)
    slides.append(s)

    # Slide 4: Part I divider
    s = make_slide(prs, color_bg=CREAM)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.4),
             "PART", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.9), Inches(12.333), Inches(2.5),
             "I", size=120, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.5), Inches(2.3), Emu(12700))
    rule.fill.solid(); rule.fill.fore_color.rgb = GOLD; rule.line.fill.background()
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.333), Inches(0.6),
             "The Vision", size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(s, 4, 18)
    slides.append(s)

    # Slide 5: 5 Commitments
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART I  ·  CORE PHILOSOPHY", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "The 5 Commitments", size=28, bold=True, color=NAVY)
    add_rule(s, 1.85)
    commitments = [
        "1.  Trading is a business",
        "2.  Discipline beats strategy",
        "3.  Consistency beats excitement",
        "4.  Patience is an edge",
        "5.  Trading should fit into life",
    ]
    add_bullets(s, Inches(1.0), Inches(2.5), Inches(11), Inches(3), commitments, size=20, color=NAVY)
    add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5),
             "These are not negotiable. The philosophy is the asset. "
             "The Master Blueprint is the philosophy expressed in commercial, product, brand, operational, and legal form.",
             size=12, color=DARK, italic=True)
    add_footer(s, 5, 18)
    slides.append(s)

    # Slide 6: Flywheel
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART II  ·  THE BUSINESS", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Value Creation Flywheel", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "flywheel", Inches(3.5), Inches(2.0), Inches(6.5), Inches(4.5))
    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
             "The flywheel is powered by philosophy, not by growth hacks.", size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 6, 18)
    slides.append(s)

    # Slide 7: Revenue Architecture
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART II  ·  REVENUE", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Revenue Architecture — 4 Tiers", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "revenue-pyramid", Inches(2.5), Inches(2.0), Inches(8.5), Inches(4.5))
    add_footer(s, 7, 18)
    slides.append(s)

    # Slide 8: Commercial Strategy
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART II  ·  COMMERCIAL STRATEGY", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "7 Strategic Commitments", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    commits = [
        "Two-tier v1: Monthly + Annual only",
        "Closed Founding cohort of 50 Members",
        "Single Trading Credit — four 'nevers'",
        "Writing-led growth engine — no paid acquisition",
        "₹1.2–1.5 Cr Founder reserve Y1, ₹15–30 Cr raise Y2",
        "Realistic Member targets: 30–50 Y1, 150–200 Y2, 350–500 Y3",
        "Two-stage launch: Design Partner (8–15) → Founding (50)",
    ]
    add_bullets(s, Inches(1.0), Inches(2.4), Inches(11), Inches(4.5), commits, size=14, color=NAVY)
    add_footer(s, 8, 18)
    slides.append(s)

    # Slide 9: Product — 12-stage lifecycle
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART III  ·  PRODUCT", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "12-Stage Member Lifecycle", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "lifecycle", Inches(0.6), Inches(2.3), Inches(12.133), Inches(3.0))
    add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.0),
             "Discovery → Integration → Leadership. The first 90 days is the identity-formation phase. "
             "The cutover from Founder-moderated to Host-moderated to Mentor-moderated is the brand's most-defensible transition discipline.",
             size=12, color=DARK, italic=True)
    add_footer(s, 9, 18)
    slides.append(s)

    # Slide 10: KPI Dashboard
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART VII  ·  FOUNDER DASHBOARD", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Top 8 KPIs", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "kpi-dashboard", Inches(0.6), Inches(2.2), Inches(12.133), Inches(3.5))
    add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.6),
             "Three most load-bearing: MRR (weekly), 12-month retention (quarterly — the gating test), "
             "Founder's writing hours/week (weekly — the growth constraint).",
             size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 10, 18)
    slides.append(s)

    # Slide 11: 24-Month Roadmap
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART VI  ·  EXECUTION", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "24-Month Strategic Roadmap", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "roadmap", Inches(0.6), Inches(2.2), Inches(12.133), Inches(3.8))
    add_footer(s, 11, 18)
    slides.append(s)

    # Slide 12: Growth Curve
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART II  ·  GROWTH", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Membership Growth — 3 Scenarios", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "growth-curve", Inches(1.5), Inches(2.2), Inches(10), Inches(4.0))
    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
             "Expected: 30–50 Y1, 150–200 Y2, 350–500 Y3. The 75% retention gate gates the Y2 capital raise.",
             size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, 12, 18)
    slides.append(s)

    # Slide 13: Risk Heatmap
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART VI  ·  RISK", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Risk Heatmap — Top 10 Risks", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "risk-heatmap", Inches(0.6), Inches(2.2), Inches(8.5), Inches(4.5))
    risks = [
        "R1: Founder burnout (Q4)",
        "R2: Retention <75%",
        "R3: Y2 raise fails",
        "R4: SEBI enforcement",
        "R5: Host departure",
    ]
    add_bullets(s, Inches(9.5), Inches(2.7), Inches(3.5), Inches(4), risks, size=12, color=NAVY)
    add_footer(s, 13, 18)
    slides.append(s)

    # Slide 14: Org Chart
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART V  ·  TEAM", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Team Structure — Cap of 8 in 24 Months", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    add_diagram(s, "org-chart", Inches(1.5), Inches(2.2), Inches(10.3), Inches(4.5))
    add_footer(s, 14, 18)
    slides.append(s)

    # Slide 15: Immediate Priorities
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART VIII  ·  BOARD", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Immediate Priorities — Next 7 Days", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    items = [
        "Set Annual price (OPEN-006)",
        "Set Monthly price (OPEN-007)",
        "Set Founding cohort cap at 50 (OPEN-008)",
        "Set Headcount safety factor at 1.5x (OPEN-009)",
        "Confirm ₹1.2–1.5 Cr Founder-capital reserve",
        "Sign external counsel retainer",
        "Adopt this Blueprint as operating reference",
    ]
    add_bullets(s, Inches(1.0), Inches(2.4), Inches(11), Inches(4.5), items, size=15, color=NAVY)
    add_footer(s, 15, 18)
    slides.append(s)

    # Slide 16: Things NOT to do
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART VIII  ·  THE NEGATION", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "10 Things NOT to Do", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    notdos = [
        "No Premium service in Y1",
        "No Dubai before retention ≥75% and 30+ applicants",
        "No regional chapter before 200+ Members in the metro",
        "No hiring above 8 in 24 months",
        "No paid acquisition — ever",
        "No discounting any tier",
        "No crossing the 'what we are NOT' line",
        "No Board override of Brand Strategist's veto",
        "No Y2 raise with growth targets, board seat, or liquidation preference above 1x",
        "No Founder weekly hours exceeding 55",
    ]
    add_bullets(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.7), notdos, size=12, color=NAVY)
    add_footer(s, 16, 18)
    slides.append(s)

    # Slide 17: Critical Risks
    s = make_slide(prs)
    add_text(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "PART VIII  ·  CRITICAL RISKS", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
             "Top 5 Critical Risks", size=24, bold=True, color=NAVY)
    add_rule(s, 1.85)
    risks_top = [
        "1. Founder burnout by Q4",
        "2. 12-month Founding retention below 75%",
        "3. Year 2 capital raise fails on philosophy-defending terms",
        "4. SEBI enforcement on Floor content",
        "5. Host departure in Year 1 or 2",
    ]
    add_bullets(s, Inches(1.0), Inches(2.4), Inches(11), Inches(3.5), risks_top, size=15, color=NAVY)
    add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5),
             "Reviewed weekly by the Founder, monthly by the Operations Excellence Consultant, "
             "and quarterly by the Leadership Team. Mitigations are operationalised.",
             size=12, color=DARK, italic=True)
    add_footer(s, 17, 18)
    slides.append(s)

    # Slide 18: Closing — "The plan begins"
    s = make_slide(prs, color_bg=NAVY)
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.333), Inches(0.4),
             "FROM THE EXECUTIVE STRATEGY OFFICE", size=10, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(2.6), Inches(1.3), Emu(12700))
    rule.fill.solid(); rule.fill.fore_color.rgb = GOLD; rule.line.fill.background()
    add_text(s, Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.2),
             "The institution is the philosophy", size=28, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8),
             "in operational form.", size=28, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
    rule2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(5.0), Inches(1.3), Emu(12700))
    rule2.fill.solid(); rule2.fill.fore_color.rgb = GOLD; rule2.line.fill.background()
    add_text(s, Inches(0.5), Inches(5.4), Inches(12.333), Inches(0.6),
             "The plan begins.", size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
             "FX TELLER  ·  STRATEGIC BLUEPRINT  ·  V1.0  ·  JULY 2026",
             size=9, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.3),
             "CONFIDENTIAL", size=8, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    # No footer on closing
    slides.append(s)

    # Save
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    print(f"  Saved: {PPTX_PATH}")
    print(f"  Slides: {len(slides)}")
    print(f"  Size: {os.path.getsize(PPTX_PATH)//1024} KB")


if __name__ == "__main__":
    main()
